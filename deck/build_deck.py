"""
Builds the presentation deck (deck/Climate_Intelligence_Challenge_Deck.pptx) using real
numbers from model/metrics.json. Run after model/train_model.py.

    py deck/build_deck.py
"""
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(__file__)
ROOT = os.path.dirname(HERE)

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
BLUE = RGBColor(0x15, 0x65, 0xC0)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)

metrics_path = os.path.join(ROOT, "model", "metrics.json")
metrics = None
if os.path.exists(metrics_path):
    with open(metrics_path) as f:
        metrics = json.load(f)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=NAVY,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def header(slide, title, subtitle=None):
    add_text(slide, 0.6, 0.35, 11, 0.8, title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, 0.6, 1.0, 11, 0.5, subtitle, size=16, color=BLUE)
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.5), Inches(11.0), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def metric_card(slide, left, top, width, value, label, color=BLUE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(34)
    r1.font.bold = True
    r1.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.color.rgb = GREY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----
    s = blank_slide(prs)
    bg(s, NAVY)
    add_text(s, 0.8, 2.4, 11.7, 1.2, "Forecasting Dangerous Heat", size=44, bold=True, color=WHITE)
    add_text(s, 0.8, 3.4, 11.7, 0.8,
             "10-Day Wet-Bulb Temperature Forecasting for Uttar Pradesh", size=22, color=RGBColor(0x9E, 0xC9, 0xFF))
    add_text(s, 0.8, 4.3, 11.7, 0.5,
             "Live demo: ather-hackathon.streamlit.app", size=17, bold=True, color=RGBColor(0x66, 0xBB, 0x6A))
    add_text(s, 0.8, 6.3, 11.7, 0.5,
             "IIIT Lucknow Climate Intelligence Challenge 2026  ·  Aether AI/ML Club", size=15, color=GREY)

    # ---- Slide 2: Problem ----
    s = blank_slide(prs)
    header(s, "The Problem: Heat Is Not Just Temperature")
    add_bullets(s, 0.6, 1.9, 6.5, 4.5, [
        "Uttar Pradesh: 240M+ people, warming faster than almost any region on Earth.",
        "Wet-Bulb Temperature (WBT) = temperature + humidity → the body's real cooling limit.",
        "Above 32°C WBT: outdoor activity should stop.",
        "Above 35°C WBT: unsurvivable beyond ~6 hours, even in shade, at rest.",
        "Current forecasting systems degrade sharply beyond 2–3 days for heat-humidity variables.",
        "Goal: push accurate WBT forecasting out to a 10-day horizon.",
    ])
    tbl_left = 7.4
    rows = [("< 28°C", "Safe", GREEN), ("28–32°C", "Caution", ORANGE),
            ("32–35°C", "Danger", RGBColor(0xE6, 0x51, 0x00)), ("> 35°C", "Fatal", RED)]
    top = 2.0
    for wbt, label, color in rows:
        box = s.shapes.add_shape(1, Inches(tbl_left), Inches(top), Inches(4.8), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{wbt}   —   {label}"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = WHITE
        top += 1.05

    # ---- Slide 3: Data ----
    s = blank_slide(prs)
    header(s, "The Data", "NASA POWER daily meteorological reanalysis")
    add_bullets(s, 0.6, 1.9, 11.5, 4.5, [
        "31 years of training data (1984–2014) across 125 districts of UP & neighboring regions.",
        "Test period 2015–2025 (dates masked as day_index); context.csv seeds the first 30-day lag window per district.",
        "Features: T2M/T2M_MAX/T2M_MIN, RH2M, QV2M, soil temp/moisture, solar radiation (all-sky & clear-sky), "
        "cloud amount, wind speed/direction, surface pressure, precipitation, evapotranspiration.",
        "WBT precomputed for the training period via the Stull (2011) approximation.",
        "Districts anonymized to relative coordinates (rel_lat, rel_lon) — no district names in the data.",
    ])

    # ---- Slide 4: Approach ----
    s = blank_slide(prs)
    header(s, "Our Approach", "Why same-day scoring beats recursive forecasting here")
    add_bullets(s, 0.6, 1.9, 11.5, 2.2, [
        "Key insight: the test period's meteorological drivers are known NASA reanalysis values, not forecasts —"
        " so we don't need to recursively predict future weather to predict future WBT.",
        "A LightGBM regressor learns WBT directly from same-day features + engineered lag/rolling statistics"
        " (1/3/7-day lags, 7-day rolling mean & std, day-over-day deltas), computed per district.",
        "The 10-day-ahead forecast for any day is simply the model's own same-day predictions on the following"
        " 10 days — this avoids compounding recursive-forecast error entirely.",
    ])
    add_text(s, 0.6, 4.4, 11.5, 0.4, "Pipeline", size=18, bold=True, color=BLUE)
    steps = ["context.csv + train.csv", "Feature engineering\n(lags, rolling stats,\npsychrometric features)",
             "LightGBM\n(same-day WBT)", "Shift predictions\n→ 10-day targets", "submission.csv /\nlive dashboard"]
    x = 0.6
    w = 2.2
    for i, step in enumerate(steps):
        box = s.shapes.add_shape(5, Inches(x), Inches(5.0), Inches(w - 0.3), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE if i % 2 == 0 else RGBColor(0x0D, 0x47, 0xA1)
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = step
        r.font.size = Pt(12)
        r.font.color.rgb = WHITE
        x += w

    # ---- Slide 5: Model performance ----
    s = blank_slide(prs)
    header(s, "Model Performance", "Validation results (held-out 15% time split)")
    if metrics:
        metric_card(s, 0.8, 2.0, 3.5, f"{metrics['wrmse_10day']:.3f}°C", "10-day Weighted RMSE", BLUE)
        metric_card(s, 4.5, 2.0, 3.5, f"{metrics['same_day_rmse']:.3f}°C", "Same-day RMSE", GREEN)
        metric_card(s, 8.2, 2.0, 3.5, str(metrics["n_features"]), "Engineered features", ORANGE)
        add_text(s, 0.8, 3.5, 11.5, 0.4, "Per-day validation RMSE (weights decay 1.0 → 0.1, Day 1 → Day 10)",
                 size=16, bold=True, color=NAVY)
        per_day = metrics["per_day_rmse"]
        chart_left, chart_top, chart_w, chart_h = 0.9, 4.0, 11.0, 2.8
        max_val = max(per_day) if per_day else 1
        bar_w = chart_w / 10
        for i, v in enumerate(per_day):
            h = (v / max_val) * (chart_h - 0.6)
            box = s.shapes.add_shape(1, Inches(chart_left + i * bar_w + 0.1), Inches(chart_top + (chart_h - 0.6) - h),
                                      Inches(bar_w - 0.2), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = BLUE
            box.line.fill.background()
            add_text(s, chart_left + i * bar_w, chart_top + chart_h - 0.55, bar_w, 0.3, f"D{i+1}",
                      size=11, align=PP_ALIGN.CENTER, color=NAVY)
            add_text(s, chart_left + i * bar_w, chart_top + (chart_h - 0.6) - h - 0.3, bar_w, 0.3, f"{v:.2f}",
                      size=10, align=PP_ALIGN.CENTER, color=GREY)
    else:
        add_text(s, 0.8, 2.5, 11, 1, "Run model/train_model.py to populate live metrics here.", size=18, color=RED)

    # ---- Slide 6: Deployment / UI ----
    s = blank_slide(prs)
    header(s, "Live Deployment", "Streamlit dashboard — interactive district risk forecasting")
    add_bullets(s, 0.6, 1.9, 11.5, 3.0, [
        "Region-wide risk map: all 125 districts, color-coded Safe/Caution/Danger/Fatal, scrub any day 2015–2025.",
        "Per-district drill-down: full 10-day WBT trajectory with danger-band overlay + peak/day-1/day-10 metrics.",
        "Model performance panel: live WRMSE, per-horizon RMSE, and feature-importance charts.",
        "Deployed on Streamlit Community Cloud directly from GitHub — no infra to manage.",
    ])
    add_text(s, 0.6, 5.2, 11.5, 0.5, "Live link: https://ather-hackathon.streamlit.app/", size=18, bold=True, color=BLUE)
    add_text(s, 0.6, 5.8, 11.5, 0.5, "GitHub: https://github.com/Laveshjadon/ATHER-HACKATHON", size=18, bold=True, color=BLUE)

    # ---- Slide 7: Impact ----
    s = blank_slide(prs)
    header(s, "Real-World Impact")
    add_bullets(s, 0.6, 1.9, 11.5, 4.5, [
        "A 10-day heat-risk horizon gives district administrations, hospitals, and outdoor-labor employers"
        " real lead time instead of a 2-3 day window.",
        "Region-wide map surfaces which districts are trending toward Danger/Fatal before it happens —"
        " usable for targeted heat-action-plan activation (water points, work-hour shifts, health advisories).",
        "Built entirely on public NASA POWER data — reproducible for any region with the same feature set.",
        "Extensible: swap in additional districts, add monsoon-onset features, or feed into a heat-action-plan"
        " alerting pipeline (SMS/WhatsApp advisories to vulnerable groups).",
    ])

    # ---- Slide 8: Next steps / close ----
    s = blank_slide(prs)
    bg(s, NAVY)
    add_text(s, 0.8, 2.6, 11.7, 1.0, "Next Steps", size=34, bold=True, color=WHITE)
    add_bullets(s, 0.8, 3.6, 11.5, 3.0, [
        "Ensemble with the LSTM/CatBoost variants already explored (weather_tcn_lstm.py) for a small accuracy gain.",
        "Add spatial neighbor-district features to capture regional heatwave propagation.",
        "Wire the dashboard into an SMS/WhatsApp alerting pipeline for on-the-ground heat advisories.",
    ], color=RGBColor(0xCF, 0xE0, 0xFF))
    add_text(s, 0.8, 6.6, 11.7, 0.5, "Thank you — Aether, IIIT Lucknow", size=16, color=GREY)

    out_path = os.path.join(HERE, "Climate_Intelligence_Challenge_Deck.pptx")
    prs.save(out_path)
    print(f"Saved {out_path}")


if __name__ == "__main__":
    build()
